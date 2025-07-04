"""
AUTO_PPT_SURAJDESIGNER.py

Automated PowerPoint Generation from Image-based Questions using OCR

This script processes a folder of scanned question images, extracts each question using Tesseract OCR, and generates a PowerPoint presentation with each question as a separate slide. Designed for educators and content creators who need to digitize and present question papers efficiently.

Author: Suraj (surajdesigner)
Website: https://surajdesigner.com
GitHub: https://github.com/designersuraj
License: MIT

For more information, visit surajdesigner.com
"""
import io
import os
import sys
import re
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
import cv2
import pytesseract

# -----------------------------
# Author: Suraj (designersuraj)
# For more info: surajdesigner.com
# -----------------------------

# --- Tesseract Path Setup (auto set, so no PATH issues) ---
# Ensure Tesseract-OCR is installed at the specified path
# This avoids PATH issues on most Windows setups
tess_path = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
if os.path.exists(tess_path):
    pytesseract.pytesseract.tesseract_cmd = tess_path
else:
    print("ERROR: Tesseract not found at", tess_path)
    print("Please check your installation.")
    sys.exit(1)

print("Python chal raha hai:", sys.executable)  # Debug: Show Python executable path
print("Python version:", sys.version)           # Debug: Show Python version

# --- File Paths ---
# Set the folder containing input images and the output PPTX file path
image_folder = r"D:\UAE_Session_2025-26\2025_New Automation Scripts\PPT Automation by Image\SUBJECT"
pptx_path = r"D:\UAE_Session_2025-26\2025_New Automation Scripts\PPT Automation by Image\PPT OUTPUT BY SUBJECT.pptx"

# --- Get all .jpg images in sorted order ---
# Only process .jpg files, sorted numerically for consistency
image_files = [f for f in os.listdir(image_folder) if f.lower().endswith('.jpg')]
def extract_number(x):
    m = re.match(r'^(\d+)', x)
    return int(m.group(1)) if m else float('inf')
image_files.sort(key=extract_number)

if not image_files:
    print(f"No .jpg images found in {image_folder}")
    sys.exit(1)

# --- Open or Create Presentation ---
# Try to open an existing PPTX, otherwise create a new one
try:
    prs = Presentation(pptx_path)
except Exception:
    prs = Presentation()

# --- Find a blank slide layout safely ---
# Prefer a blank layout for clean slides, fallback to first layout if not found
blank_slide_layout = None
for layout in prs.slide_layouts:
    if "blank" in layout.name.lower():
        blank_slide_layout = layout
        break
if blank_slide_layout is None:
    blank_slide_layout = prs.slide_layouts[0]  # fallback

# --- Process each image file ---
question_count = 0
for image_file in image_files:
    image_path = os.path.join(image_folder, image_file)
    print(f"Processing {image_path} ...")
    img = cv2.imread(image_path)
    if img is None:
        print("ERROR: Image not found at", image_path)
        continue

    # Convert to grayscale and apply threshold for better OCR
    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    _, thresh = cv2.threshold(gray, 180, 255, cv2.THRESH_BINARY_INV + cv2.THRESH_OTSU)
    ocr_img = gray  # or use thresh for more contrast

    # --- OCR Config ---
    # Use Tesseract to extract text data with bounding boxes
    custom_config = r'--oem 3 --psm 6'
    data = pytesseract.image_to_data(ocr_img, config=custom_config, output_type=pytesseract.Output.DICT)

    # Build a list of (text, top, left, width, height) for all words
    words = []
    for i in range(len(data['text'])):
        t = data['text'][i].strip()
        if t:
            words.append({
                'text': t,
                'top': data['top'][i],
                'left': data['left'][i],
                'width': data['width'][i],
                'height': data['height'][i],
                'idx': i
            })

    # Join all words into a single string (with spaces)
    full_text = ' '.join([w['text'] for w in words])

    # Find question starts (e.g., 1., 2., ...) robustly: match '^\d+\s*\.?$', and only if left is near left edge
    question_starts = []
    option4_indices = []
    img_width = img.shape[1]
    for i, w in enumerate(words):
        if re.match(r'^\d+\s*\.?$', w['text']) and w['left'] < 0.2 * img_width:
            question_starts.append(i)
        if w['text'] == '(4)':
            option4_indices.append(i)

    # For each question, crop from start to (4) (and 2 more words if next question number is not immediately after)
    question_boxes = []
    for i, start_idx in enumerate(question_starts):
        # End: next question number index or end of image
        if i < len(question_starts) - 1:
            end_idx = question_starts[i+1] - 1
        else:
            end_idx = len(words) - 1
        block_indices = list(range(start_idx, end_idx + 1))
        block_tops = [words[k]['top'] for k in block_indices]
        block_lefts = [words[k]['left'] for k in block_indices]
        block_widths = [words[k]['width'] for k in block_indices]
        block_heights = [words[k]['height'] for k in block_indices]
        top = min(block_tops)
        left = min(block_lefts)
        right = max([block_lefts[j] + block_widths[j] for j in range(len(block_lefts))])
        bottom = max([block_tops[j] + block_heights[j] for j in range(len(block_tops))]) + 40  # EXTRA_PAD
        # Clamp to image size
        left = max(0, left)
        right = min(img.shape[1], right)
        top = max(0, top)
        bottom = min(img.shape[0], bottom)
        question_boxes.append((left, top, right, bottom))

    if not question_boxes:
        print(f"No questions detected in {image_file}! Please check OCR or image quality.")
        continue

    # --- Crop and Paste Each Question to PPTX ---
    for i, box in enumerate(question_boxes):
        x1, y1, x2, y2 = box
        cropped = img[y1:y2, x1:x2]
        pil_img = Image.fromarray(cv2.cvtColor(cropped, cv2.COLOR_BGR2RGB))
        image_stream = io.BytesIO()
        pil_img.save(image_stream, format='PNG')
        image_stream.seek(0)

        # Get image size in inches for accurate placement in PPTX
        img_width, img_height = pil_img.size
        dpi = 96  # PowerPoint assumes 96 dpi
        img_width_in = img_width / dpi
        img_height_in = img_height / dpi

        # Center the image on the slide
        # Handle NoneType for slide_width/height (fallback to default 10x7.5 inches)
        slide_width = prs.slide_width if prs.slide_width is not None else 9144000  # 10 inches in EMU
        slide_height = prs.slide_height if prs.slide_height is not None else 6858000  # 7.5 inches in EMU
        slide_width_in = slide_width / 914400
        slide_height_in = slide_height / 914400
        left = Inches((slide_width_in - img_width_in) / 2)
        top = Inches((slide_height_in - img_height_in) / 2)

        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(image_stream, left, top, width=Inches(img_width_in), height=Inches(img_height_in))
        question_count += 1
        print(f"Added question {question_count} (from {image_file}) to slide.")

# Save the final presentation
prs.save(pptx_path)
print(f"Done! Total {question_count} questions added to PPTX.")

# -----------------------------
# Script by surajdesigner | surajdesigner.com
# -----------------------------
