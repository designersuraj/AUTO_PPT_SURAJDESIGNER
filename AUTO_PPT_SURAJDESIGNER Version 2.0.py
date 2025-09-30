"""
AUTO_PPT_SURAJDESIGNER.py

Automated PowerPoint Generation from Image-based Questions using OCR

This script processes a folder of scanned question images, extracts each question using Tesseract OCR, and generates a PowerPoint presentation with each question as a separate slide. Designed for educators and content creators who need to digitize and present question papers efficiently.

Author: Suraj (surajdesigner)
Website: https://surajdesigner.com
GitHub: https://github.com/designersuraj
License: MIT

For more information, visit surajdesigner.com
"""
import io
import os
import sys
import re
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
import cv2
import pytesseract

# -----------------------------
# Author: Suraj (designersuraj)
# For more info: surajdesigner.com
# -----------------------------

# --- Tesseract Path Setup (auto set, so no PATH issues) ---
# Ensure Tesseract-OCR is installed at the specified path
# This avoids PATH issues on most Windows setups
tess_path = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
if os.path.exists(tess_path):
    pytesseract.pytesseract.tesseract_cmd = tess_path
else:
    print("ERROR: Tesseract not found at", tess_path)
    print("Please check your installation.")
    sys.exit(1)

print("Python Running:", sys.executable)  # Debug: Show Python executable path
print("Python version:", sys.version)           # Debug: Show Python version

# --- File Paths ---
# Set the folder containing input images and the output PPTX file path
image_folder = r"D:\PNCF\2809 VID SOL\PHY"
pptx_path = r"D:\PNCF\2809 VID SOL\2809_PHYSICS_PNCF-10th_MINOR TEST-04_(VIDEO SOLUTION).pptx"

# --- Get all .jpg images in sorted order ---
# Only process .jpg files, sorted numerically for consistency
image_files = [f for f in os.listdir(image_folder) if f.lower().endswith('.jpg')]
def extract_number(x):
    m = re.match(r'^(\d+)', x)
    return int(m.group(1)) if m else float('inf')
image_files.sort(key=extract_number)

if not image_files:
    print(f"No .jpg images found in {image_folder}")
    sys.exit(1)

# --- Placement Settings ---
# Images will be scaled to fit within these margins and positioned per alignment settings
# Horizontal options: "center", "left", "right"
# Vertical options:   "center", "top"
horizontal_align = "left"
vertical_align = "top"
content_margin_in = 0.5  # margin from slide edges on all sides
left_guideline_in = 1.0   # used when horizontal_align == "left"
right_guideline_in = 1.0  # used when horizontal_align == "right" (distance from right edge)
top_guideline_in = 1.0    # used when vertical_align == "top"

# Prefer using an anchor box on a template slide to position/fit images
use_template_anchor = True
anchor_shape_candidates = [
    "IMAGE_BOX",
    "image_box",  # allow lowercase exact match
]
debug_anchor_scan = True  # print shapes discovered to help configure anchor

def _emu_to_inches(value_emu):
    return value_emu / 914400

def _matches_anchor(shape):
    try:
        name = (getattr(shape, "name", "") or "").strip().lower()
        alt_text = (getattr(shape, "alternative_text", "") or getattr(shape, "alt_text", "") or "").strip().lower()
        text = ""
        if hasattr(shape, "text_frame") and shape.text_frame is not None:
            text = (shape.text_frame.text or "").strip().lower()
        cands = [c.lower() for c in anchor_shape_candidates]
        tokens = [name, alt_text, text]
        # exact token match only (not substring) to avoid matching backgrounds/placeholders accidentally
        return any(token in cands for token in tokens)
    except Exception:
        return False

def _is_fullslide_shape(shape, slide_width_in, slide_height_in, threshold: float = 0.97):
    try:
        w_in = _emu_to_inches(shape.width)
        h_in = _emu_to_inches(shape.height)
        return (w_in >= slide_width_in * threshold) and (h_in >= slide_height_in * threshold)
    except Exception:
        return False

def _iter_shapes_recursive(container):
    """Yield shapes recursively (handles grouped shapes)."""
    try:
        for shape in container.shapes:
            yield shape
            if hasattr(shape, "shapes"):
                for child in _iter_shapes_recursive(shape):
                    yield child
    except Exception:
        return

def find_template_anchor_box(presentation):
    """Return (left_in, top_in, width_in, height_in, source_label, matched_name) from slide, layout, or master.
    Priority: first slide -> its layout -> any layout -> any master.
    """
    slide_width_in = (presentation.slide_width or 9144000) / 914400
    slide_height_in = (presentation.slide_height or 6858000) / 914400

    # PRIORITY 1) Masters (take reference from Slide Master if present)
    try:
        for master in presentation.slide_masters:
            if debug_anchor_scan:
                print(f"-- Scan (priority): master '{getattr(master,'name','master')}' shapes --")
                for shp in _iter_shapes_recursive(master):
                    try:
                        print(
                            f"  master shape: name='{getattr(shp,'name','')}', alt='{getattr(shp,'alternative_text','') or getattr(shp,'alt_text','')}', size=({_emu_to_inches(shp.width):.2f}x{_emu_to_inches(shp.height):.2f}) in"
                        )
                    except Exception:
                        pass
            for shape in _iter_shapes_recursive(master):
                if _matches_anchor(shape):
                    if _is_fullslide_shape(shape, slide_width_in, slide_height_in):
                        print(f"Anchor candidate on master '{getattr(master, 'name', 'master')}' skipped: spans entire slide")
                        continue
                    return (
                        _emu_to_inches(shape.left),
                        _emu_to_inches(shape.top),
                        _emu_to_inches(shape.width),
                        _emu_to_inches(shape.height),
                        f"master:{getattr(master, 'name', 'master')}",
                        getattr(shape, "name", "") or getattr(shape, "alternative_text", "") or getattr(shape, "alt_text", "") or "",
                    )
    except Exception:
        pass

    # 1) First slide
    try:
        first_slide = presentation.slides[0]
        if debug_anchor_scan:
            print("-- Scan: first slide shapes --")
            for shape in _iter_shapes_recursive(first_slide):
                try:
                    print(
                        f"  slide shape: name='{getattr(shape,'name','')}', alt='{getattr(shape,'alternative_text','') or getattr(shape,'alt_text','')}', size=({_emu_to_inches(shape.width):.2f}x{_emu_to_inches(shape.height):.2f}) in"
                    )
                except Exception:
                    pass
        for shape in _iter_shapes_recursive(first_slide):
            if _matches_anchor(shape):
                if _is_fullslide_shape(shape, slide_width_in, slide_height_in):
                    print("Anchor candidate on slide_0 skipped: spans entire slide")
                    continue
                return (
                    _emu_to_inches(shape.left),
                    _emu_to_inches(shape.top),
                    _emu_to_inches(shape.width),
                    _emu_to_inches(shape.height),
                    "slide_0",
                    getattr(shape, "name", "") or getattr(shape, "alternative_text", "") or getattr(shape, "alt_text", "") or "",
                )
    except Exception:
        pass

    # 2) First slide's layout
    try:
        layout = presentation.slides[0].slide_layout
        if debug_anchor_scan:
            print(f"-- Scan: layout '{layout.name}' shapes --")
            for shape in _iter_shapes_recursive(layout):
                try:
                    print(
                        f"  layout shape: name='{getattr(shape,'name','')}', alt='{getattr(shape,'alternative_text','') or getattr(shape,'alt_text','')}', size=({_emu_to_inches(shape.width):.2f}x{_emu_to_inches(shape.height):.2f}) in"
                    )
                except Exception:
                    pass
        for shape in _iter_shapes_recursive(layout):
            if _matches_anchor(shape):
                if _is_fullslide_shape(shape, slide_width_in, slide_height_in):
                    print(f"Anchor candidate on layout '{layout.name}' skipped: spans entire slide")
                    continue
                return (
                    _emu_to_inches(shape.left),
                    _emu_to_inches(shape.top),
                    _emu_to_inches(shape.width),
                    _emu_to_inches(shape.height),
                    f"layout:{layout.name}",
                    getattr(shape, "name", "") or getattr(shape, "alternative_text", "") or getattr(shape, "alt_text", "") or "",
                )
    except Exception:
        pass

    # 3) Any layout
    try:
        for layout in presentation.slide_layouts:
            if debug_anchor_scan:
                print(f"-- Scan: layout '{layout.name}' shapes (all layouts) --")
                for shp in _iter_shapes_recursive(layout):
                    try:
                        print(
                            f"  layout shape: name='{getattr(shp,'name','')}', alt='{getattr(shp,'alternative_text','') or getattr(shp,'alt_text','')}', size=({_emu_to_inches(shp.width):.2f}x{_emu_to_inches(shp.height):.2f}) in"
                        )
                    except Exception:
                        pass
            for shape in _iter_shapes_recursive(layout):
                if _matches_anchor(shape):
                    if _is_fullslide_shape(shape, slide_width_in, slide_height_in):
                        print(f"Anchor candidate on layout '{layout.name}' skipped: spans entire slide")
                        continue
                    return (
                        _emu_to_inches(shape.left),
                        _emu_to_inches(shape.top),
                        _emu_to_inches(shape.width),
                        _emu_to_inches(shape.height),
                        f"layout:{layout.name}",
                        getattr(shape, "name", "") or getattr(shape, "alternative_text", "") or getattr(shape, "alt_text", "") or "",
                    )
    except Exception:
        pass

    # 4) Masters
    try:
        for master in presentation.slide_masters:
            if debug_anchor_scan:
                print(f"-- Scan: master '{getattr(master,'name','master')}' shapes --")
                for shp in _iter_shapes_recursive(master):
                    try:
                        print(
                            f"  master shape: name='{getattr(shp,'name','')}', alt='{getattr(shp,'alternative_text','') or getattr(shp,'alt_text','')}', size=({_emu_to_inches(shp.width):.2f}x{_emu_to_inches(shp.height):.2f}) in"
                        )
                    except Exception:
                        pass
            for shape in _iter_shapes_recursive(master):
                if _matches_anchor(shape):
                    if _is_fullslide_shape(shape, slide_width_in, slide_height_in):
                        print(f"Anchor candidate on master '{getattr(master, 'name', 'master')}' skipped: spans entire slide")
                        continue
                    return (
                        _emu_to_inches(shape.left),
                        _emu_to_inches(shape.top),
                        _emu_to_inches(shape.width),
                        _emu_to_inches(shape.height),
                        f"master:{getattr(master, 'name', 'master')}",
                        getattr(shape, "name", "") or getattr(shape, "alternative_text", "") or getattr(shape, "alt_text", "") or "",
                    )
    except Exception:
        pass

    # 5) Fallback: choose the largest 'Rectangle*' shape on any layout/master that is not full-slide
    try:
        best = None
        best_source = None
        # scan layouts
        for layout in presentation.slide_layouts:
            for shape in _iter_shapes_recursive(layout):
                try:
                    name = (getattr(shape, "name", "") or "").lower()
                    if name.startswith("rectangle") and not _is_fullslide_shape(shape, slide_width_in, slide_height_in):
                        area = _emu_to_inches(shape.width) * _emu_to_inches(shape.height)
                        if best is None or area > best[0]:
                            best = (area, shape)
                            best_source = f"layout:{layout.name}"
                except Exception:
                    continue
        # scan masters
        for master in presentation.slide_masters:
            for shape in _iter_shapes_recursive(master):
                try:
                    name = (getattr(shape, "name", "") or "").lower()
                    if name.startswith("rectangle") and not _is_fullslide_shape(shape, slide_width_in, slide_height_in):
                        area = _emu_to_inches(shape.width) * _emu_to_inches(shape.height)
                        if best is None or area > best[0]:
                            best = (area, shape)
                            best_source = f"master:{getattr(master, 'name', 'master')}"
                except Exception:
                    continue
        if best is not None:
            shape = best[1]
            print(f"Fallback anchor chosen from {best_source} (name='{getattr(shape,'name','')}')")
            return (
                _emu_to_inches(shape.left),
                _emu_to_inches(shape.top),
                _emu_to_inches(shape.width),
                _emu_to_inches(shape.height),
                best_source,
                getattr(shape, "name", ""),
            )
    except Exception:
        pass

    return None

# --- Open or Create Presentation ---
# Try to open an existing PPTX, otherwise create a new one
try:
    prs = Presentation(pptx_path)
except Exception:
    prs = Presentation()

# Try to read an anchor box from the first slide of the existing presentation
template_anchor_result = find_template_anchor_box(prs) if use_template_anchor else None
template_anchor_box_in = None
if template_anchor_result is not None:
    anchor_left_in, anchor_top_in, anchor_w_in, anchor_h_in, anchor_source, anchor_name = template_anchor_result
    template_anchor_box_in = (anchor_left_in, anchor_top_in, anchor_w_in, anchor_h_in)
    print(f"Anchor found from {anchor_source} (shape='{anchor_name}'): left={anchor_left_in:.2f} in, top={anchor_top_in:.2f} in, width={anchor_w_in:.2f} in, height={anchor_h_in:.2f} in")
else:
    print("No anchor shape found; using margin/alignment settings.")

# --- Find a blank slide layout safely ---
# Prefer a blank layout for clean slides, fallback to first layout if not found
blank_slide_layout = None
for layout in prs.slide_layouts:
    if "blank" in layout.name.lower():
        blank_slide_layout = layout
        break
if blank_slide_layout is None:
    blank_slide_layout = prs.slide_layouts[0]  # fallback

# --- Process each image file ---
question_count = 0
for image_file in image_files:
    image_path = os.path.join(image_folder, image_file)
    print(f"Processing {image_path} ...")
    img = cv2.imread(image_path)
    if img is None:
        print("ERROR: Image not found at", image_path)
        continue

    # Convert to grayscale and apply threshold for better OCR
    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    _, thresh = cv2.threshold(gray, 180, 255, cv2.THRESH_BINARY_INV + cv2.THRESH_OTSU)
    ocr_img = gray  # or use thresh for more contrast

    # --- OCR Config ---
    # Use Tesseract to extract text data with bounding boxes
    custom_config = r'--oem 3 --psm 6'
    data = pytesseract.image_to_data(ocr_img, config=custom_config, output_type=pytesseract.Output.DICT)

    # Build a list of (text, top, left, width, height) for all words
    words = []
    for i in range(len(data['text'])):
        t = data['text'][i].strip()
        if t:
            words.append({
                'text': t,
                'top': data['top'][i],
                'left': data['left'][i],
                'width': data['width'][i],
                'height': data['height'][i],
                'idx': i
            })

    # Join all words into a single string (with spaces)
    full_text = ' '.join([w['text'] for w in words])

    # Find question starts (e.g., 1., 2., ...) robustly: match '^\d+\s*\.?$', and only if left is near left edge
    question_starts = []
    option4_indices = []
    img_width = img.shape[1]
    for i, w in enumerate(words):
        if re.match(r'^\d+\s*\.?$', w['text']) and w['left'] < 0.2 * img_width:
            question_starts.append(i)
        if w['text'] == '(4)':
            option4_indices.append(i)

    # For each question, crop from start to (4) (and 2 more words if next question number is not immediately after)
    question_boxes = []
    for i, start_idx in enumerate(question_starts):
        # End: next question number index or end of image
        if i < len(question_starts) - 1:
            end_idx = question_starts[i+1] - 1
        else:
            end_idx = len(words) - 1
        block_indices = list(range(start_idx, end_idx + 1))
        block_tops = [words[k]['top'] for k in block_indices]
        block_lefts = [words[k]['left'] for k in block_indices]
        block_widths = [words[k]['width'] for k in block_indices]
        block_heights = [words[k]['height'] for k in block_indices]
        top = min(block_tops)
        left = min(block_lefts)
        right = max([block_lefts[j] + block_widths[j] for j in range(len(block_lefts))])
        bottom = max([block_tops[j] + block_heights[j] for j in range(len(block_tops))]) + 40  # EXTRA_PAD
        # Clamp to image size
        left = max(0, left)
        right = min(img.shape[1], right)
        top = max(0, top)
        bottom = min(img.shape[0], bottom)
        question_boxes.append((left, top, right, bottom))

    if not question_boxes:
        print(f"No questions detected in {image_file}! Please check OCR or image quality.")
        continue

    # --- Crop and Paste Each Question to PPTX ---
    for i, box in enumerate(question_boxes):
        x1, y1, x2, y2 = box
        cropped = img[y1:y2, x1:x2]
        pil_img = Image.fromarray(cv2.cvtColor(cropped, cv2.COLOR_BGR2RGB))
        image_stream = io.BytesIO()
        pil_img.save(image_stream, format='PNG')
        image_stream.seek(0)

        # Get image size in inches for accurate placement in PPTX
        img_width, img_height = pil_img.size
        dpi = 96  # PowerPoint assumes 96 dpi
        img_width_in = img_width / dpi
        img_height_in = img_height / dpi

        # Handle NoneType for slide_width/height (fallback to default 10x7.5 inches)
        slide_width = prs.slide_width if prs.slide_width is not None else 9144000  # 10 inches in EMU
        slide_height = prs.slide_height if prs.slide_height is not None else 6858000  # 7.5 inches in EMU
        slide_width_in = slide_width / 914400
        slide_height_in = slide_height / 914400

        # Compute maximum drawable area considering margins or template anchor
        if template_anchor_box_in is not None:
            anchor_left_in, anchor_top_in, anchor_w_in, anchor_h_in = template_anchor_box_in
            max_w_in = max(0.1, anchor_w_in)
            max_h_in = max(0.1, anchor_h_in)
        else:
            max_w_in = max(0.1, slide_width_in - 2 * content_margin_in)
            max_h_in = max(0.1, slide_height_in - 2 * content_margin_in)

        # Scale image to fit within the drawable area without exceeding slide bounds
        scale = min(max_w_in / img_width_in, max_h_in / img_height_in, 1.0)
        draw_w_in = img_width_in * scale
        draw_h_in = img_height_in * scale

        # Position based on alignment settings and guidelines
        ha = horizontal_align.lower()
        va = vertical_align.lower()

        if template_anchor_box_in is not None:
            # Place the image inside the anchor box according to horizontal_align/vertical_align
            if ha == "left":
                left_in = anchor_left_in
            elif ha == "right":
                left_in = anchor_left_in + max_w_in - draw_w_in
            else:  # center
                left_in = anchor_left_in + (max_w_in - draw_w_in) / 2

            if va == "top":
                top_in = anchor_top_in
            else:  # center
                top_in = anchor_top_in + (max_h_in - draw_h_in) / 2

            # Clamp within anchor bounds (just in case)
            left_in = max(anchor_left_in, min(left_in, anchor_left_in + max_w_in - draw_w_in))
            top_in = max(anchor_top_in, min(top_in, anchor_top_in + max_h_in - draw_h_in))

            # Debug print for each slide placement
            print(f"Placing image inside anchor: x={left_in:.2f} in, y={top_in:.2f} in, w={draw_w_in:.2f} in, h={draw_h_in:.2f} in")
        else:
            if ha == "left":
                left_in = max(content_margin_in, left_guideline_in)
                right_limit_in = slide_width_in - content_margin_in
                if left_in + draw_w_in > right_limit_in:
                    left_in = max(content_margin_in, right_limit_in - draw_w_in)
            elif ha == "right":
                right_edge_target_in = max(content_margin_in, right_guideline_in)
                left_in = slide_width_in - right_edge_target_in - draw_w_in
                if left_in < content_margin_in:
                    left_in = content_margin_in
            else:
                left_in = (slide_width_in - draw_w_in) / 2

            if va == "top":
                top_in = max(content_margin_in, top_guideline_in)
                bottom_limit_in = slide_height_in - content_margin_in
                if top_in + draw_h_in > bottom_limit_in:
                    top_in = max(content_margin_in, bottom_limit_in - draw_h_in)
            else:
                top_in = (slide_height_in - draw_h_in) / 2

        left = Inches(left_in)
        top = Inches(top_in)

        # If we discovered an anchor on a specific layout, try to use that same layout for new slides
        if template_anchor_result is not None and anchor_source.startswith("layout:"):
            # Find layout by name; fallback to blank
            layout_name = anchor_source.split(":", 1)[1]
            chosen_layout = blank_slide_layout
            try:
                for layout in prs.slide_layouts:
                    if getattr(layout, "name", "") == layout_name:
                        chosen_layout = layout
                        break
            except Exception:
                pass
            slide = prs.slides.add_slide(chosen_layout)
        else:
            slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(image_stream, left, top, width=Inches(draw_w_in), height=Inches(draw_h_in))
        question_count += 1
        print(f"Added question {question_count} (from {image_file}) to slide.")

# Save the final presentation
prs.save(pptx_path)
print(f"Done! Total {question_count} questions added to PPTX.")

# -----------------------------
# Script by surajdesigner | surajdesigner.com
# -----------------------------
